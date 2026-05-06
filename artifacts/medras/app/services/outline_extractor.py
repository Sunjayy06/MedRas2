"""Outline-step text extraction.

Wraps the existing extractor in `plagiarism_analyzer` and adds PowerPoint
(`.pptx`) support. All other formats (`.pdf`, `.docx`, `.txt`, `.md`) and all
error semantics are inherited unchanged.
"""

from __future__ import annotations

from io import BytesIO
from typing import Set

from app.services import plagiarism_analyzer as _pa

# Re-export for convenience so callers can `except outline_extractor.UploadExtractionError`.
UploadExtractionError = _pa.UploadExtractionError
UnsupportedFileError = _pa.UnsupportedFileError
PasswordProtectedError = _pa.PasswordProtectedError
ImageOnlyPdfError = _pa.ImageOnlyPdfError
DocumentTooLargeError = _pa.DocumentTooLargeError
CorruptedFileError = _pa.CorruptedFileError

SUPPORTED_EXTENSIONS: Set[str] = {".pdf", ".docx", ".pptx", ".txt", ".md"}


def _extract_pptx(content: bytes) -> str:
    """Extract every readable string from a .pptx file using python-pptx.

    We deliberately walk *every shape on every slide* (recursing into group
    shapes), pulling text from:
      * title / body placeholders and free-floating text boxes
      * table cells
      * chart titles, axis titles, series names, and category labels
      * SmartArt / freeform shapes that expose a text frame
      * speaker notes
      * slide-layout and slide-master placeholders that act as headers,
        footers, page numbers, dates, and watermark text
    Anything that raises is silently skipped — we never want one malformed
    shape to kill the whole extraction.
    """
    try:
        from pptx import Presentation  # type: ignore
        from pptx.exc import PackageNotFoundError  # type: ignore
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
    except ImportError as exc:  # pragma: no cover — dep is pinned
        raise CorruptedFileError(
            "PowerPoint support is not installed on this server."
        ) from exc

    try:
        prs = Presentation(BytesIO(content))
    except PackageNotFoundError as exc:
        raise CorruptedFileError(
            "We could not open this PowerPoint file. It may be password-protected, "
            "an older .ppt format, or not actually a .pptx file. Please save it as a "
            "modern .pptx and try again."
        ) from exc
    except Exception as exc:
        raise CorruptedFileError(
            "We could not read this PowerPoint file. Please re-export it as .pptx and try again."
        ) from exc

    def _text_from_frame(tf) -> list:
        out = []
        try:
            for para in tf.paragraphs:
                line = "".join((run.text or "") for run in para.runs).strip()
                if line:
                    out.append(line)
        except Exception:
            pass
        return out

    def _text_from_chart(chart) -> list:
        out = []
        try:
            if getattr(chart, "has_title", False) and chart.has_title:
                t = (chart.chart_title.text_frame.text or "").strip()
                if t: out.append("[Chart title] " + t)
        except Exception: pass
        # Axis titles
        for axis_attr in ("category_axis", "value_axis"):
            try:
                ax = getattr(chart, axis_attr, None)
                if ax is not None and getattr(ax, "has_title", False) and ax.has_title:
                    t = (ax.axis_title.text_frame.text or "").strip()
                    if t: out.append(f"[{axis_attr.replace('_', ' ').title()}] " + t)
            except Exception: pass
        # Series names + category labels
        try:
            for series in chart.series:
                nm = (getattr(series, "name", "") or "").strip()
                if nm: out.append("[Series] " + nm)
        except Exception: pass
        try:
            cats = list(chart.plots[0].categories)  # type: ignore[index]
            cleaned = [str(c).strip() for c in cats if str(c or "").strip()]
            if cleaned:
                out.append("[Chart categories] " + ", ".join(cleaned))
        except Exception: pass
        return out

    def _walk_shapes(shapes) -> list:
        out = []
        for shape in shapes:
            try:
                # Recurse into groups first
                if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                    try:
                        out.extend(_walk_shapes(shape.shapes))
                    except Exception:
                        pass
                    continue
                # Text frames (titles, body, text boxes, freeforms with text)
                if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                    out.extend(_text_from_frame(shape.text_frame))
                # Tables
                if getattr(shape, "has_table", False) and shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            cell_text = (cell.text or "").strip()
                            if cell_text:
                                out.append(cell_text)
                # Charts
                if getattr(shape, "has_chart", False) and shape.has_chart:
                    out.extend(_text_from_chart(shape.chart))
            except Exception:
                continue
        return out

    parts = []
    for idx, slide in enumerate(prs.slides, 1):
        slide_text_parts: list = []

        # Slide layout placeholders — usually carry header/footer/date/page-number
        # and any text the deck author put on the master template.
        try:
            for ph in slide.slide_layout.placeholders:
                if getattr(ph, "has_text_frame", False) and ph.has_text_frame:
                    for line in _text_from_frame(ph.text_frame):
                        slide_text_parts.append("[Layout] " + line)
        except Exception:
            pass

        # Every shape on the slide itself
        try:
            slide_text_parts.extend(_walk_shapes(slide.shapes))
        except Exception:
            pass

        # Speaker notes
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
                if notes:
                    slide_text_parts.append("[Notes] " + notes)
        except Exception:
            pass

        # De-duplicate consecutive identical lines (layout placeholders often
        # repeat the title verbatim).
        cleaned: list = []
        prev = None
        for line in slide_text_parts:
            if line and line != prev:
                cleaned.append(line)
                prev = line

        if cleaned:
            parts.append(f"--- Slide {idx} ---\n" + "\n".join(cleaned))

    return "\n\n".join(parts)


def extract_text(filename: str, content: bytes) -> str:
    """Extract text from any supported upload type for the outline step.

    Adds .pptx handling on top of the existing plagiarism extractor. Other
    extensions and error semantics are unchanged. May raise the same
    ``UploadExtractionError`` subclasses as ``plagiarism_analyzer``.
    """
    name = (filename or "").lower().strip()
    if name.endswith(".pptx"):
        return _extract_pptx(content)
    if name.endswith(".ppt"):
        raise UnsupportedFileError(
            "Old-format .ppt files aren't supported. Please open the file in PowerPoint "
            "and 'Save As' a .pptx, then try again."
        )
    # Delegate everything else to the existing extractor (.pdf/.docx/.txt/.md).
    return _pa.extract_text_from_upload(filename, content)
